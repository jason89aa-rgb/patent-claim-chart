"""PPTX Claim Chart Export — 무효성 검토자료 에디토리얼 디자인.

디자인 시스템 (사용자 제공 "특허 클레임 차트 디자인 개선" 덱 기준):
- 캔버스 20 x 11.25 inch, 배경 #F2F2F3
- 각진 사각형 + 0.8pt 검정 테두리 + 헤어라인 구분선 (둥근 모서리 없음)
- 아이브로우(작은 파란 라벨) + 큰 제목의 인쇄물 스타일 헤더
- 판단은 신호등 색 텍스트 + 옅은 칩 배경
- 청구항(좌) / 선행문헌(우) 세로 분할, 도면은 흰 카드 + 다크 FIG 배지
"""
import os
import re

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

from core.project import ProjectData, CaseInfo, Claim, MappingEntry
from core.region_capture import (capture_for_mappings,
                                 group_mappings_by_region,
                                 build_region_index)
from utils.term_format import split_by_terms, evidence_chunks
from utils.errlog import log_exception

# ------------------------------------------------------------ 캔버스
SLIDE_W = Inches(20)
SLIDE_H = Inches(11.25)

HEAD_FONT = "Barlow Condensed"     # 제목·라벨 (없으면 PowerPoint가 대체)
BODY_FONT = "Barlow"               # 본문
EA_FONT = "맑은 고딕"               # 한글 폰트 (동아시아 폴백)

# ------------------------------------------------------------ 색상 토큰
BG      = RGBColor(0xF2, 0xF2, 0xF3)   # 슬라이드 배경
BAND    = RGBColor(0xE9, 0xE9, 0xEA)   # 헤더/코멘트 밴드
LABEL_BG = RGBColor(0xED, 0xED, 0xEE)  # 라벨 셀
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
LINE    = RGBColor(0x1F, 0x1F, 0x1F)   # 테두리·구분선
INK     = RGBColor(0x1D, 0x1F, 0x20)   # 본문 텍스트
MUTED   = RGBColor(0x7A, 0x7A, 0x7D)   # 보조 텍스트
ACCENT  = RGBColor(0x41, 0x61, 0x80)   # 아이브로우·요소 ID
ACCENT2 = RGBColor(0x59, 0x80, 0xA6)   # 번호
DEEP    = RGBColor(0x2C, 0x45, 0x5D)   # 진한 라벨
BADGE   = RGBColor(0x1D, 0x2D, 0x3D)   # FIG 배지 배경

# 판단: (글자색, 칩 배경)
JUDGMENT_STYLE = {
    "일치":    (RGBColor(0x1B, 0x68, 0x3E), RGBColor(0xD7, 0xF4, 0xE0)),
    "부분일치": (RGBColor(0x8A, 0x60, 0x00), RGBColor(0xFB, 0xEE, 0xC9)),
    "불일치":  (RGBColor(0x9D, 0x35, 0x33), RGBColor(0xF8, 0xE3, 0xE1)),
    "미판단":  (RGBColor(0x7A, 0x7A, 0x7D), RGBColor(0xE9, 0xE9, 0xEA)),
}

INTERPRETATION_LABELS = {
    "문언침해": "Lit.",
    "균등론":   "DOE",
    "넓게해석": "Broad",
    "좁게해석": "Narrow",
}

FIGURE_REF_NOTE = "(도면 표시 참조)"


def _doc_label(path: str) -> str:
    """문서 표기: 확장자(.pdf)와 복사본 접미사를 떼고 등록번호만 남긴다."""
    if not path:
        return "선행문헌"
    name = os.path.splitext(os.path.basename(path))[0]
    name = re.sub(r"\s*\(\d+\)\s*$", "", name)        # "US11367770 (1)"
    name = re.sub(r"\s*-\s*복사본\s*$", "", name)        # "... - 복사본"
    name = re.sub(r"_[0-9a-f]{12}$", "", name)         # 붙여넣은 텍스트 문서
    return name.strip() or "선행문헌"


def _docs_label(paths: list) -> str:
    """여러 선행문헌을 한 줄로 (청구항 한 장에 모두 실을 때)."""
    names = [_doc_label(p) for p in paths if p]
    if not names:
        return "선행문헌"
    if len(names) <= 2:
        return " · ".join(names)
    return f"{names[0]} 외 {len(names) - 1}건"


# ------------------------------------------------------------ 기본 유틸
def _no_shadow(shape):
    try:
        shape.shadow.inherit = False
    except Exception:
        pass


def _rgb(t) -> RGBColor:
    return RGBColor(*tuple(t))


def _readable(rgb) -> RGBColor:
    """밝은 배경(#F2F2F3)에서 읽히도록 너무 밝은 색은 어둡게 보정.

    색상(hue)은 유지하므로 앱 화면의 요소 색과 대응 관계가 유지된다.
    """
    r, g, b = tuple(rgb)
    lum = 0.299 * r + 0.587 * g + 0.114 * b
    if lum > 150:                      # 밝으면 눌러서 대비 확보
        k = 150.0 / max(lum, 1.0)
        r, g, b = int(r * k), int(g * k), int(b * k)
    return RGBColor(r, g, b)


def _rect(slide, left, top, width, height,
          fill=None, line=None, line_w=0.8):
    """각진 사각형 (이 디자인은 둥근 모서리를 쓰지 않는다)."""
    sp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                Inches(left), Inches(top),
                                Inches(width), Inches(height))
    if fill is None:
        sp.fill.background()
    else:
        sp.fill.solid()
        sp.fill.fore_color.rgb = fill
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line
        sp.line.width = Pt(line_w)
    _no_shadow(sp)
    sp.text_frame.word_wrap = True
    return sp


def _hairline(slide, left, top, width):
    """0.01인치 실선 구분선 (표 테두리 대신 사용)."""
    return _rect(slide, left, top, width, 0.01, fill=LINE)


def _vline(slide, left, top, height):
    return _rect(slide, left, top, 0.01, height, fill=LINE)


def _apply_font(run, size, bold=False, color=INK, condensed=False):
    f = run.font
    f.size = Pt(size)
    f.bold = bold
    f.color.rgb = color
    f.name = HEAD_FONT if condensed else BODY_FONT     # latin
    # 한글은 동아시아 폰트로 (Barlow에는 한글 글리프가 없음)
    rPr = run._r.get_or_add_rPr()
    if rPr.find(qn("a:ea")) is None:
        ea = rPr.makeelement(qn("a:ea"), {})
        ea.set("typeface", EA_FONT)
        rPr.append(ea)


def _tf(shape, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, pad=0.0):
    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Pt(pad)
    tf.margin_top = tf.margin_bottom = Pt(0)
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    return p


def _text(slide, left, top, width, height, text,
          size=18, bold=False, color=INK, condensed=False,
          align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, wrap=True):
    """텍스트 박스 (배경·테두리 없음)."""
    box = slide.shapes.add_textbox(Inches(left), Inches(top),
                                   Inches(width), Inches(height))
    p = _tf(box, align, anchor)
    box.text_frame.word_wrap = wrap
    run = p.add_run()
    run.text = text
    _apply_font(run, size, bold, color, condensed)
    return box


def _rich(slide, left, top, width, height, chunks,
          size=18, base=INK, align=PP_ALIGN.LEFT,
          anchor=MSO_ANCHOR.TOP):
    """[(텍스트, 색상 또는 None), ...] 을 한 문단에 이어 붙인다.

    색상이 있는 조각(매칭 요소)은 굵게 — 청구항과 선행문헌에서
    같은 요소가 같은 색으로 보이게 하는 핵심.
    """
    box = slide.shapes.add_textbox(Inches(left), Inches(top),
                                   Inches(width), Inches(height))
    p = _tf(box, align, anchor)
    for chunk, color in chunks:
        if not chunk:
            continue
        run = p.add_run()
        run.text = chunk
        if color:
            _apply_font(run, size, True, _readable(color))
        else:
            _apply_font(run, size, False, base)
    return box


def _text_w(text: str, size: float, condensed: bool = True) -> float:
    """텍스트 렌더 폭(인치) 추정.

    한글은 정사각 글리프(1em ≈ size pt)이고 Barlow Condensed 라틴은
    약 0.38em로 훨씬 좁다 — 둘을 따로 계산해야 칩 안에서 한글이
    세로로 깨지지 않는다.
    """
    ko = sum(1 for c in text if ord(c) > 0x2E80)
    latin = len(text) - ko
    lat_f = 0.0060 if condensed else 0.0068
    return ko * size * 0.0150 + latin * size * lat_f


def _chip(slide, left, top, text, fg, bg, size=18, pad=0.11):
    """옅은 배경 + 진한 글자의 상태 칩. 텍스트 길이에 맞춰 폭 자동."""
    width = max(0.70, _text_w(text, size) + pad * 2)
    _rect(slide, left, top, width, 0.43, fill=bg)
    _text(slide, left + pad, top + 0.02, width - pad * 2, 0.40,
          text, size=size, bold=True, color=fg, condensed=True,
          anchor=MSO_ANCHOR.MIDDLE, wrap=False)
    return width


def _dot(slide, left, top, color, hollow=False):
    """매트릭스용 작은 사각 상태 점."""
    if hollow:
        _rect(slide, left, top, 0.17, 0.17, fill=None, line=color)
    else:
        _rect(slide, left, top, 0.17, 0.17, fill=color)


def _judgment(judgment: str, interp: str = "") -> tuple:
    """(라벨, 글자색, 칩배경)"""
    fg, bg = JUDGMENT_STYLE.get(judgment, JUDGMENT_STYLE["미판단"])
    label = judgment
    il = INTERPRETATION_LABELS.get(interp, "")
    if il:
        label = f"{judgment}·{il}"
    return label, fg, bg


def _slide(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg = s.background.fill
    bg.solid()
    bg.fore_color.rgb = BG
    return s


def _header(slide, eyebrow: str, title: str):
    """아이브로우 + 큰 제목 (모든 본문 슬라이드 공통)."""
    _text(slide, 0.83, 0.52, 12.0, 0.43, eyebrow,
          size=18, bold=True, color=ACCENT, condensed=True)
    _text(slide, 0.83, 1.04, 18.4, 0.70, title,
          size=39, bold=True, color=INK, condensed=True)


def _color_maps(data: ProjectData) -> dict:
    colors = {}
    for claim in data.claims:
        for elem in claim.elements:
            colors[elem.element_id] = tuple(elem.color_rgb)
    for t in data.terms:
        colors[t.term_id] = tuple(t.color_rgb)
    return colors


def _fig_card(slide, left, top, height, img_path, label="",
              min_w=2.6, max_w=8.4) -> float:
    """흰 카드 + 가운데 정렬 도면 + 다크 배지.

    카드 폭을 도면 종횡비에 맞춰 정하므로 도면이 카드 안에서
    작게 떠 보이지 않는다. 사용한 폭을 반환.
    """
    inset = 0.10
    iw = ih = None
    if img_path and os.path.exists(img_path):
        try:
            from PIL import Image
            with Image.open(img_path) as im:
                iw, ih = im.size
        except Exception as e:
            print(f"[export_pptx] image read error: {e}")

    if iw and ih:
        width = (iw / ih) * (height - inset * 2) + inset * 2
        width = max(min_w, min(max_w, width))
    else:
        width = min_w

    _rect(slide, left, top, width, height, fill=WHITE, line=LINE)
    if iw and ih:
        try:
            scale = min((width - inset * 2) / iw, (height - inset * 2) / ih)
            w, h = iw * scale, ih * scale
            slide.shapes.add_picture(
                img_path,
                Inches(left + (width - w) / 2),
                Inches(top + (height - h) / 2),
                width=Inches(w), height=Inches(h))
        except Exception as e:
            print(f"[export_pptx] image insert error: {e}")

    if label:
        bw = _text_w(label, 18) + 0.26
        bl = left + width - bw - 0.10
        _rect(slide, bl, top + 0.11, bw, 0.45, fill=BADGE)
        _text(slide, bl + 0.13, top + 0.14, bw - 0.26, 0.40, label,
              size=18, bold=True, color=WHITE, condensed=True,
              anchor=MSO_ANCHOR.MIDDLE, wrap=False)
    return width


# ------------------------------------------------------------ 표지
def _cover_slide(prs, case_info: CaseInfo, title: str):
    s = _slide(prs)
    _rect(s, 0.83, 1.15, 18.33, 8.96, fill=None, line=LINE)

    _text(s, 14.5, 0.52, 4.7, 0.43, "무효성 검토자료 · 내부용",
          size=18, bold=True, color=ACCENT, condensed=True,
          align=PP_ALIGN.RIGHT)
    _text(s, 0.83, 0.52, 8.0, 0.43, "CLAIM CHART",
          size=18, bold=True, color=ACCENT, condensed=True)

    _text(s, 1.60, 2.30, 17.0, 1.70, title or "특허 Claim Chart",
          size=72, bold=True, color=ACCENT, condensed=True)

    ident = (case_info.registration_number or
             case_info.application_number or "")
    if ident:
        _text(s, 1.60, 4.35, 17.0, 0.40,
              "대상 특허 등록번호" if case_info.registration_number
              else "대상 특허 출원번호",
              size=22, bold=True, color=INK, condensed=True)
        _text(s, 1.60, 4.80, 17.0, 1.30, ident,
              size=66, bold=True, color=INK, condensed=True)

    rows = [
        ("발명의 명칭", case_info.title),
        ("출원인", case_info.applicant),
        ("우선일 · 출원일", " · ".join(
            x for x in (case_info.priority_date,
                        case_info.application_date) if x)),
        ("등록일", case_info.registration_date),
        ("패밀리", ", ".join(case_info.family_patents)),
    ]
    y = 6.60
    for label, value in rows:
        if not value:
            continue
        _hairline(s, 1.60, y, 16.8)
        _text(s, 1.60, y + 0.18, 3.2, 0.45, label,
              size=18, bold=True, color=DEEP, condensed=True)
        _text(s, 5.00, y + 0.18, 13.4, 0.45, value,
              size=18.75, color=INK)
        y += 0.72
        if y > 9.60:
            break


# ------------------------------------------------------------ 서지정보
def _caseinfo_slide(prs, data: ProjectData):
    ci = data.case_info
    rows = [
        ("발명의 명칭", ci.title),
        ("출원번호", ci.application_number),
        ("등록번호", ci.registration_number),
        ("우선일", ci.priority_date),
        ("출원일", ci.application_date),
        ("등록일", ci.registration_date),
        ("출원인", ci.applicant),
        ("패밀리 특허", ", ".join(ci.family_patents)),
    ]
    rows = [(k, v) for k, v in rows if v]
    if not rows and not data.doc_paths:
        return

    s = _slide(prs)
    _header(s, "서지정보", "대상 특허 · 선행문헌")

    # ── 좌측: 대상 특허
    card_h = min(7.0, 0.70 + 0.66 * max(len(rows), 1) + 0.10)
    _rect(s, 0.83, 2.04, 8.96, card_h, fill=None, line=LINE)
    _rect(s, 0.84, 2.05, 8.94, 0.70, fill=BAND)
    _text(s, 1.09, 2.23, 6.0, 0.40, "대상 특허",
          size=19.5, bold=True, color=INK, condensed=True)
    _hairline(s, 0.84, 2.74, 8.94)

    y = 2.75
    for label, value in rows:
        _rect(s, 0.84, y, 2.08, 0.65, fill=LABEL_BG)
        _text(s, 1.07, y + 0.13, 1.85, 0.45, label,
              size=18, bold=True, color=DEEP, condensed=True)
        _text(s, 3.15, y + 0.13, 6.55, 0.45, value,
              size=18.75, color=INK)
        y += 0.66
        _hairline(s, 0.84, y - 0.01, 8.94)

    # ── 우측: 선행문헌 목록
    docs = data.doc_paths or list(dict.fromkeys(
        m.doc_path for m in data.mappings if m.doc_path))
    doc_h = min(7.0, 0.70 + 0.66 * max(len(docs), 1) + 0.10)
    _rect(s, 10.23, 2.04, 8.96, doc_h, fill=None, line=LINE)
    _rect(s, 10.24, 2.05, 8.94, 0.70, fill=BAND)
    _text(s, 10.49, 2.23, 6.0, 0.40, "선행문헌",
          size=19.5, bold=True, color=INK, condensed=True)
    _hairline(s, 10.24, 2.74, 8.94)

    y = 2.75
    for i, dp in enumerate(docs, 1):
        n_map = sum(1 for m in data.mappings if m.doc_path == dp)
        _rect(s, 10.24, y, 2.08, 0.65, fill=LABEL_BG)
        _text(s, 10.47, y + 0.13, 1.85, 0.45, f"D{i}",
              size=18, bold=True, color=DEEP, condensed=True)
        _text(s, 12.55, y + 0.13, 5.0, 0.45, _doc_label(dp),
              size=18.75, color=INK)
        _text(s, 17.60, y + 0.13, 1.5, 0.45, f"매핑 {n_map}건",
              size=18, color=MUTED, align=PP_ALIGN.RIGHT)
        y += 0.66
        _hairline(s, 10.24, y - 0.01, 8.94)
    if not docs:
        _text(s, 10.49, 2.95, 8.0, 0.45, "(등록된 선행문헌 없음)",
              size=18.75, color=MUTED)


# ------------------------------------------------------------ 요소 범례
def _term_legend_slide(prs, terms: list):
    if not terms:
        return
    s = _slide(prs)
    _header(s, "매칭 요소", "요소 색상 범례")
    _text(s, 0.83, 1.80, 18.4, 0.40,
          "같은 색 = 청구항의 요소와 선행문헌의 대응부분이 서로 대응함을 의미합니다.",
          size=18.75, color=MUTED)

    col_w, per_col = 9.10, 12
    for i, t in enumerate(terms):
        col, row = i // per_col, i % per_col
        x = 0.83 + col * (col_w + 0.30)
        y = 2.45 + row * 0.62
        if x + col_w > 20.0:
            break
        _hairline(s, x, y, col_w)
        rgb = _readable(t.color_rgb)
        _rect(s, x, y + 0.18, 0.22, 0.22, fill=rgb)
        _text(s, x + 0.42, y + 0.13, 1.2, 0.40, t.term_id,
              size=20.25, bold=True, color=rgb, condensed=True)
        _text(s, x + 1.60, y + 0.14, col_w - 1.7, 0.40, t.text,
              size=19.5, color=INK)


# ------------------------------------------------------------ Type A
# 아이브로우를 없애고 카드를 위로 올려 본문 높이를 최대한 확보
# (청구항 하나가 되도록 한 장에 들어가도록)
_A_CARD = (0.83, 0.62, 18.33, 10.00)    # left, top, w, h
_A_SPLIT = 6.04                          # 청구항 | 선행문헌 세로 분할선
_A_L = 1.07                              # 좌측 본문 x
_A_LW = 4.89                             # 좌측 본문 폭
_A_R = 6.28                              # 우측 본문 x
_A_RW = 12.63                            # 우측 본문 폭 (카드 안쪽 여백 유지)
_A_BODY_TOP = 1.39                       # 카드 상단 + 헤더밴드(0.77)
_A_BODY_BOT = 10.40
_FIG_H = 3.20                            # 도면 카드 높이
_FIG_GAP = 0.16


def _a_header(prs, claim: Claim, doc_label: str, page: int,
              verdict: tuple | None, card_h: float | None = None):
    s = _slide(prs)
    suffix = f"  ({page})" if page > 1 else ""
    l, t, w, h = _A_CARD
    if card_h is not None:                 # 내용 높이에 맞춰 카드 축소
        h = max(2.2, min(h, card_h))
    _rect(s, l, t, w, h, fill=BG, line=LINE)
    _rect(s, l + 0.01, t + 0.01, w - 0.02, 0.76, fill=BAND)
    _text(s, 1.09, t + 0.15, 4.9, 0.51,
          f"청구항 {claim.claim_number}"
          + ("" if claim.is_independent
             else f" (제{claim.parent_claim}항 종속)") + suffix,
          size=18, bold=True, color=INK, condensed=True)
    _text(s, _A_R, t + 0.23, 8.0, 0.34, f"선행문헌 · {doc_label}",
          size=18, bold=True, color=INK, condensed=True)
    if verdict:
        label, fg, bg = verdict
        cw = max(1.6, _text_w(label, 18) + 0.34)
        _rect(s, l + w - cw - 0.25, t + 0.15, cw, 0.47, fill=bg)
        _text(s, l + w - cw - 0.11, t + 0.19, cw - 0.28, 0.43, label,
              size=18, bold=True, color=fg, condensed=True,
              anchor=MSO_ANCHOR.MIDDLE, wrap=False)
    _hairline(s, l + 0.01, _A_BODY_TOP, w - 0.02)
    return s, t + h          # 슬라이드, 카드 하단 y


def _wrap_h(text: str, width_in: float, size: float,
            line_h: float = None) -> float:
    """텍스트 렌더 높이 추정.

    실측 기준(18pt 본문): 라틴 약 0.145in/자, 줄 간격 약 0.26in.
    단어 단위 줄바꿈 여유(6%)를 두고 넉넉히 잡아야 '분석 의견' 같은
    아래 블록이 카드 밖으로 잘리지 않는다.
    """
    if line_h is None:
        line_h = size * 0.0160
    if not text:
        return line_h
    ko = sum(1 for c in text if ord(c) > 0x2E80)
    latin = len(text) - ko
    w = ko * size * 0.0139 + latin * size * 0.0080
    usable = max(width_in * 0.94, 0.5)
    lines = max(1, int(w / usable) + 1)
    return lines * line_h


def _a_elements(claim: Claim, mappings: list) -> list:
    """좌측에 쌓을 구성요소 항목 (요소 ID + 판단 + 문언)."""
    by_elem: dict = {}
    for m in mappings:
        by_elem.setdefault(m.element_id, []).append(m)

    items = []
    for elem in claim.elements:
        # 요소 번호(1A…)와 판단 칩은 표시하지 않는다 — 청구항 문언에
        # 지면을 최대한 내주기 위함 (판단은 카드 헤더의 종합 배지로 표시)
        items.append({
            "elem": elem,
            "h": _wrap_h(elem.text, _A_LW, 18) + 0.26,
        })
    return items


# 명세서 문단으로 볼 최소 길이. 도면 캡션/라벨("VDDL power line in Fig. 7")은
# 이보다 짧아 도면 섹션에 남고, 긴 문단은 '명세서 대응 부분'으로만 간다.
_PARA_MIN = 60


def _is_paragraph(text: str) -> bool:
    return len((text or "").strip()) >= _PARA_MIN


def _a_evidence(mappings: list) -> list:
    """우측에 실을 명세서 대응 문단 (중복 제거).

    도면 라벨 수준의 짧은 텍스트(VDDL 등)는 제외 — 그 단어는 도면
    이미지 안에 색으로 표시되므로 문단으로 다시 적지 않는다.
    """
    seen, out = set(), []
    for m in mappings:
        t = (m.extracted_text or "").strip()
        if not _is_paragraph(t) or t in seen:
            continue
        seen.add(t)
        out.append(m)
    return out


def _a_figure_groups(mappings: list) -> list:
    """도면 섹션에 넣을 영역만 추린다.

    명세서 문단을 드래그한 영역(긴 텍스트가 있는 매핑)은 제외 — 같은
    내용이 '명세서 대응 부분'에 그대로 실리므로 캡처까지 넣으면 중복이다.
    """
    out = []
    for members in group_mappings_by_region(mappings).values():
        if not any(_is_paragraph(m.extracted_text) for m in members):
            out.append(members)
    return out


def _a_right_blocks(groups: list, paras: list, note: str) -> list:
    """우측 컬럼을 블록 단위로 (페이지 분할이 가능하도록)."""
    blocks = []
    if groups:
        blocks.append({"kind": "figs", "groups": groups,
                       "h": _FIG_H + 0.30})
    for m in paras:
        blocks.append({
            "kind": "para", "m": m,
            "h": _wrap_h(m.extracted_text or "", _A_RW, 18) + 0.16,
        })
    blocks.append({
        "kind": "note", "text": note,
        "h": 0.40 + max(0.60, _wrap_h(note, _A_RW, 18)) + 0.10,
    })
    return blocks


def _paginate(items: list, avail: float, extra_fn=None) -> list:
    """높이 기준으로 항목을 페이지에 나눈다."""
    pages, cur, used = [], [], 0.0
    for it in items:
        extra = extra_fn(it, cur) if extra_fn else 0.0
        if cur and used + it["h"] + extra > avail:
            pages.append(cur)
            cur, used = [], 0.0
            extra = extra_fn(it, cur) if extra_fn else 0.0
        cur.append(it)
        used += it["h"] + extra
    if cur:
        pages.append(cur)
    return pages or [[]]


def _right_extra(block: dict, cur: list) -> float:
    """페이지 첫 블록이 문단이면 '명세서 대응 부분' 라벨 자리를 확보."""
    return 0.40 if (not cur and block["kind"] == "para") else 0.0


def _type_a_slide(prs, claim: Claim, mappings: list, doc_label: str,
                  terms: list, colors: dict, verdict=None,
                  note: str = ""):
    """청구항(좌) × 선행문헌(우) 대비.

    우측은 구성요소별로 칸을 나누지 않고, 도면을 맨 위에 모아 배치한 뒤
    그 아래에 명세서 대응 문단과 분석 의견이 들어간다.
    좌/우 각각 넘치면 페이지를 늘린다 (내용이 잘리지 않도록).
    """
    avail = _A_BODY_BOT - (_A_BODY_TOP + 0.18)
    left_pages = _paginate(_a_elements(claim, mappings), avail)
    right_pages = _paginate(
        _a_right_blocks(_a_figure_groups(mappings),
                        _a_evidence(mappings), note),
        avail, _right_extra)

    for pno in range(max(len(left_pages), len(right_pages))):
        page_items = left_pages[pno] if pno < len(left_pages) else []
        page_blocks = right_pages[pno] if pno < len(right_pages) else []

        left_h = sum(it["h"] for it in page_items)
        right_h = sum(b["h"] for b in page_blocks) + \
            (0.40 if page_blocks and page_blocks[0]["kind"] == "para" else 0)
        card_h = 0.77 + max(left_h, right_h) + 0.24

        s, _bottom = _a_header(prs, claim, doc_label, pno + 1, verdict,
                               card_h=card_h)
        _vline(s, _A_SPLIT, _A_BODY_TOP, max(left_h, right_h) + 0.24)

        # ── 좌: 구성요소 (칸 분할 없이 세로로 쌓기)
        y = _A_BODY_TOP + 0.18
        for it in page_items:
            _rich(s, _A_L, y, _A_LW, it["h"] - 0.26,
                  split_by_terms(it["elem"].text, terms), size=18)
            y += it["h"]
        if not page_items and pno == 0:
            _text(s, _A_L, _A_BODY_TOP + 0.45, 4.6, 0.5,
                  "구성요소가 없습니다.", size=18.75, color=MUTED)

        # ── 우: 도면 → 명세서 문단 → 분석 의견
        _a_draw_right(s, page_blocks, terms, colors)


def _a_draw_right(s, blocks: list, terms: list, colors: dict):
    """우측 블록들을 순서대로 그린다."""
    y = _A_BODY_TOP + 0.18
    para_label_done = False

    for i, b in enumerate(blocks):
        if b["kind"] == "figs":
            x = _A_R
            for members in b["groups"]:
                img = capture_for_mappings(members, colors, terms=terms)
                w = _fig_card(s, x, y, _FIG_H, img,
                              f"p.{members[0].page + 1}",
                              min_w=2.2, max_w=6.0)
                x += w + _FIG_GAP
                if x > _A_R + _A_RW - 2.2:
                    break
            y += b["h"]

        elif b["kind"] == "para":
            if not para_label_done:
                _hairline(s, _A_R, y - 0.14, _A_RW)
                label = "명세서 대응 부분"
                if i == 0:      # 페이지가 넘어와 이어지는 경우
                    label += " (계속)"
                _text(s, _A_R, y, 5.0, 0.34, label,
                      size=18, bold=True, color=DEEP, condensed=True)
                y += 0.40
                para_label_done = True
            _rich(s, _A_R, y, _A_RW, b["h"] - 0.16,
                  evidence_chunks(b["m"], terms), size=18)
            y += b["h"]

        else:   # note
            _hairline(s, _A_R, y - 0.14, _A_RW)
            _text(s, _A_R, y, 4.0, 0.34, "분석 의견",
                  size=18, bold=True, color=DEEP, condensed=True)
            y += 0.40
            if b["text"]:
                _text(s, _A_R, y, _A_RW,
                      max(0.60, _wrap_h(b["text"], _A_RW, 18)),
                      b["text"], size=18, color=INK)
            y += b["h"] - 0.40


# ------------------------------------------------------------ Type B
def _type_b_slide(prs, claim: Claim, mappings: list, doc_label: str,
                  terms: list, colors: dict):
    """구성요소 1개당 슬라이드 1장 — 도면을 크게."""
    by_elem: dict = {}
    for m in mappings:
        by_elem.setdefault(m.element_id, []).append(m)
    region_index = build_region_index(mappings)

    for elem in claim.elements:
        elem_maps = by_elem.get(elem.element_id, [])
        s = _slide(prs)
        _text(s, 0.83, 0.52, 14.0, 0.43,
              f"청구항 {claim.claim_number} · {doc_label} 대비",
              size=18, bold=True, color=ACCENT, condensed=True)
        _text(s, 0.83, 1.04, 18.4, 0.70,
              f"구성요소 {elem.element_id}",
              size=39, bold=True, color=INK, condensed=True)

        if elem_maps:
            lab, fg, bg = _judgment(elem_maps[0].judgment,
                                    elem_maps[0].interpretation)
        else:
            lab, fg, bg = _judgment("미판단")
        cw = max(1.8, _text_w(lab, 19.5) + 0.36)
        _rect(s, 19.19 - cw, 1.10, cw, 0.52, fill=bg)
        _text(s, 19.19 - cw + 0.14, 1.15, cw - 0.28, 0.45, lab,
              size=19.5, bold=True, color=fg, condensed=True,
              anchor=MSO_ANCHOR.MIDDLE, wrap=False)

        # 좌: 청구항 문언
        _rect(s, 0.83, 2.10, 8.20, 7.40, fill=None, line=LINE)
        _rect(s, 0.84, 2.11, 8.18, 0.60, fill=BAND)
        _text(s, 1.09, 2.24, 6.0, 0.40, "청구항 기재",
              size=18, bold=True, color=DEEP, condensed=True)
        _rich(s, 1.09, 2.95, 7.70, 5.60,
              split_by_terms(elem.text, terms), size=21)

        notes = "\n".join(m.note for m in elem_maps if m.note)
        if notes:
            _hairline(s, 0.84, 8.55, 8.18)
            _text(s, 1.09, 8.70, 7.70, 0.70, f"비고: {notes}",
                  size=18, color=MUTED)

        # 우: 선행문헌 대응
        _rect(s, 9.40, 2.10, 9.79, 7.40, fill=None, line=LINE)
        _rect(s, 9.41, 2.11, 9.77, 0.60, fill=BAND)
        _text(s, 9.66, 2.24, 8.0, 0.40, f"선행문헌 · {doc_label}",
              size=18, bold=True, color=DEEP, condensed=True)

        regions = list(group_mappings_by_region(elem_maps).values())
        if not regions:
            _text(s, 9.66, 5.40, 9.3, 0.50, "(매핑된 선행문헌 없음)",
                  size=19.5, color=MUTED, align=PP_ALIGN.CENTER)
            continue

        y = 2.90
        for region_maps in regions[:2]:
            m0 = region_maps[0]
            _gk, gmembers = region_index.get(m0.mapping_id,
                                             (None, region_maps))
            _text(s, 9.66, y, 9.3, 0.36,
                  f"{_doc_label(m0.doc_path)} · p.{m0.page + 1}",
                  size=18, bold=True, color=ACCENT, condensed=True)
            y += 0.42

            img = capture_for_mappings(gmembers, colors, terms=terms)
            _fig_card(s, 9.66, y, 2.55, img, f"p.{m0.page + 1}",
                      min_w=3.0, max_w=9.28)
            y += 2.70

            if m0.extracted_text:
                _rich(s, 9.66, y, 9.28, 0.80,
                      evidence_chunks(m0, terms), size=18)
                y += 0.95


# ------------------------------------------------------------ Type C
def _type_c_slide(prs, claim: Claim, all_mappings: list, doc_paths: list,
                  terms: list):
    """구성요소 × 선행문헌 대응 매트릭스."""
    s = _slide(prs)
    _header(s, "대응 현황 요약",
            f"청구항 {claim.claim_number} · 구성요소 × 선행문헌 대응표")

    labels = [_doc_label(p) for p in doc_paths] or ["선행문헌"]
    n = len(labels)
    n_rows = max(len(claim.elements), 1)
    # 본문은 2.90에서 시작(헤더 0.67 + 여백 0.15)하므로 그만큼 더해야
    # 마지막 행이 잘리지 않는다.
    card_h = min(8.3, 0.82 + 0.86 * n_rows + 0.10)

    _rect(s, 0.83, 2.08, 18.33, card_h, fill=None, line=LINE)
    _rect(s, 0.84, 2.09, 18.32, 0.67, fill=BAND)

    elem_x, text_x, text_w = 1.11, 2.47, 8.89
    doc_x0, doc_w = 11.62, min(2.40, 4.6 / max(n, 1) + 1.0)
    verdict_x = 16.41

    _text(s, elem_x, 2.23, 1.0, 0.43, "요소",
          size=18, bold=True, color=INK, condensed=True)
    _text(s, text_x, 2.23, text_w, 0.43, "청구항 내용",
          size=18, bold=True, color=INK, condensed=True)
    for i, lab in enumerate(labels[:2]):
        _text(s, doc_x0 + i * doc_w, 2.23, doc_w - 0.2, 0.43,
              f"D{i+1}", size=18, bold=True, color=INK, condensed=True)
    _text(s, verdict_x, 2.23, 2.6, 0.43, "판단",
          size=18, bold=True, color=INK, condensed=True)
    _hairline(s, 0.84, 2.76, 18.32)

    y = 2.90
    for elem in claim.elements:
        if y + 0.86 > 2.08 + card_h:
            break
        _text(s, elem_x, y + 0.06, 1.0, 0.49, elem.element_id,
              size=21, bold=True, color=ACCENT, condensed=True)
        summary = elem.text if len(elem.text) <= 68 \
            else elem.text[:68].rstrip() + "…"
        _rich(s, text_x, y + 0.08, text_w, 0.46,
              split_by_terms(summary, terms), size=19.5)

        first = None
        for i, dp in enumerate(doc_paths[:2]):
            cell = [m for m in all_mappings
                    if m.element_id == elem.element_id and m.doc_path == dp]
            x = doc_x0 + i * doc_w
            if cell:
                lab, fg, _bg = _judgment(cell[0].judgment)
                _dot(s, x, y + 0.22, fg)
                _text(s, x + 0.27, y + 0.14, doc_w - 0.35, 0.36,
                      "개시" if cell[0].judgment == "일치" else lab,
                      size=18.75, bold=True, color=fg, condensed=True)
                first = first or cell[0]
            else:
                _dot(s, x, y + 0.22, MUTED, hollow=True)
                _text(s, x + 0.27, y + 0.14, 0.6, 0.36, "—",
                      size=18.75, bold=True, color=MUTED, condensed=True)

        if first:
            lab, fg, _bg = _judgment(first.judgment, first.interpretation)
        else:
            lab, fg, _bg = _judgment("미판단")
        _dot(s, verdict_x, y + 0.22, fg,
             hollow=(fg == JUDGMENT_STYLE["미판단"][0]))
        _text(s, verdict_x + 0.27, y + 0.14, 2.4, 0.36, lab,
              size=18.75, bold=True, color=fg, condensed=True)

        y += 0.86
        _hairline(s, 0.84, y - 0.10, 18.32)


# ------------------------------------------------------------ 종합
def _summary_slide(prs, data: ProjectData):
    claims = [c for c in data.claims if c.elements]
    if not claims:
        return
    s = _slide(prs)
    _header(s, "종합", "청구항별 대응 요약")

    col_w, gap = 8.94, 0.46
    for i, claim in enumerate(claims[:4]):
        col, row = i % 2, i // 2
        x = 0.83 + col * (col_w + gap)
        y = 2.13 + row * 2.55
        maps = [m for m in data.mappings
                if m.claim_number == claim.claim_number]
        mapped = {m.element_id for m in maps}
        total = len(claim.elements)
        done = len(mapped & {e.element_id for e in claim.elements})
        judgments = [m.judgment for m in maps]
        if judgments and all(j == "일치" for j in judgments) and done == total:
            verdict, fg = "전 구성요소 대응", JUDGMENT_STYLE["일치"][0]
        elif any(j == "불일치" for j in judgments):
            verdict, fg = "불일치 요소 있음", JUDGMENT_STYLE["불일치"][0]
        elif done < total:
            verdict, fg = f"미대응 {total - done}개", MUTED
        else:
            verdict, fg = "검토 중", JUDGMENT_STYLE["부분일치"][0]

        _rect(s, x, y, col_w, 2.16, fill=None, line=LINE)
        box = s.shapes.add_textbox(Inches(x + 0.39), Inches(y + 0.34),
                                   Inches(col_w - 0.7), Inches(0.45))
        p = _tf(box)
        r1 = p.add_run()
        r1.text = f"청구항 {claim.claim_number}  "
        _apply_font(r1, 25.5, True, INK, condensed=True)
        r2 = p.add_run()
        r2.text = verdict
        _apply_font(r2, 19.5, True, fg, condensed=True)

        _text(s, x + 0.39, y + 0.95, col_w - 0.78, 0.95,
              f"구성요소 {total}개 중 {done}개 대응 · 매핑 {len(maps)}건",
              size=19.5, color=INK)

    if data.case_info.notes:
        _rect(s, 0.83, 7.60, 18.33, 1.05, fill=BAND, line=LINE)
        _text(s, 1.09, 7.75, 1.6, 0.43, "담당자 의견",
              size=18, bold=True, color=DEEP, condensed=True)
        _text(s, 2.75, 7.75, 16.2, 0.80, data.case_info.notes,
              size=18, color=INK)


# ------------------------------------------------------------ 진입점
def export_pptx(data: ProjectData, output_path: str,
                template_type: str = "A",
                include_cover: bool = True):
    """PPTX 내보내기. 성공 시 None, 실패 시 오류 메시지(str) 반환."""
    try:
        prs = Presentation()
        prs.slide_width = SLIDE_W
        prs.slide_height = SLIDE_H

        terms = data.terms
        colors = _color_maps(data)

        if include_cover:
            _cover_slide(prs, data.case_info,
                         data.case_info.title or "무효성 검토")
            _caseinfo_slide(prs, data)
        if terms:
            _term_legend_slide(prs, terms)

        for claim in data.claims:
            claim_maps = [m for m in data.mappings
                          if m.claim_number == claim.claim_number]
            doc_paths = list(dict.fromkeys(
                m.doc_path for m in claim_maps))

            if template_type == "C":
                _type_c_slide(prs, claim, claim_maps, doc_paths, terms)
                continue

            if template_type == "B":
                for dp in (doc_paths or [""]):
                    dmaps = [m for m in claim_maps if m.doc_path == dp]
                    _type_b_slide(prs, claim, dmaps, _doc_label(dp),
                                  terms, colors)
                continue

            # Type A: 청구항 하나당 한 장 (선행문헌이 여러 건이어도 합쳐서)
            js = [m.judgment for m in claim_maps]
            if js and all(j == "일치" for j in js):
                verdict = _judgment("일치")
            elif any(j == "불일치" for j in js):
                verdict = _judgment("불일치")
            elif js:
                verdict = _judgment("부분일치")
            else:
                verdict = _judgment("미판단")
            note = " / ".join(
                dict.fromkeys(m.note for m in claim_maps if m.note))
            _type_a_slide(prs, claim, claim_maps, _docs_label(doc_paths),
                          terms, colors, verdict=verdict, note=note)

        if template_type != "C":
            _summary_slide(prs, data)

        prs.save(output_path)
        return None
    except Exception as e:
        log_exception("PPTX 내보내기")
        return f"{type(e).__name__}: {e}"
